"""
文档解析层

统一文档解析入口，支持 PDF / DOCX / MD / HTML / TXT / CSV 等格式。
自动识别文件类型，调用对应解析器，提取文本和表格内容。
表格自动转换为 Markdown 格式，保持结构化信息。
"""

from __future__ import annotations

import asyncio
import io
import logging
import os
import re
from typing import Any, Optional

import aiohttp

from .models import Document, compute_hash, now_iso, uuid

logger = logging.getLogger(__name__)


class DocumentParser:
    """统一文档解析器

    支持多种文件格式，自动识别类型并调用对应解析方法。
    表格统一转为 Markdown 格式，便于后续切片和检索。

    Usage::

        parser = DocumentParser()
        doc = await parser.parse("/path/to/file.pdf", tenant_id="tenant_001")
    """

    SUPPORTED_FORMATS: dict[str, str] = {
        "pdf":  "parse_pdf",
        "docx": "parse_docx",
        "md":   "parse_markdown",
        "html": "parse_html",
        "txt":  "parse_text",
        "csv":  "parse_csv",
        "xlsx": "parse_xlsx",
        "pptx": "parse_pptx",
    }

    # OCR 配置
    OCR_SERVICE_URL: str = "http://ocr-service:8003/ocr"

    def __init__(self, ocr_service_url: Optional[str] = None) -> None:
        """初始化解析器

        Args:
            ocr_service_url: OCR 微服务地址（用于扫描型 PDF）
        """
        if ocr_service_url:
            self.OCR_SERVICE_URL = ocr_service_url

    # ──────────────────────────── 公共入口 ────────────────────────────

    async def parse(self, file_path: str, tenant_id: str) -> Document:
        """解析文档

        Args:
            file_path: 文件路径
            tenant_id: 租户 ID

        Returns:
            Document 对象

        Raises:
            ValueError: 不支持的文件格式
            FileNotFoundError: 文件不存在
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        # 1. 识别格式
        ext: str = file_path.rsplit(".", 1)[-1].lower()
        if ext not in self.SUPPORTED_FORMATS:
            raise ValueError(f"不支持的文件格式: {ext}")

        # 2. 读取文件内容
        raw: bytes = await self._read_file(file_path)

        # 3. 选择解析器并执行
        parser_method_name = self.SUPPORTED_FORMATS[ext]
        parser_method = getattr(self, parser_method_name)
        sections: list[str] = await parser_method(raw)

        # 4. 计算文件哈希
        file_hash: str = compute_hash(raw)

        # 5. 提取元数据
        metadata: dict[str, Any] = self._extract_metadata(raw, ext, file_path)

        return Document(
            doc_id=uuid(),
            tenant_id=tenant_id,
            source_name=os.path.basename(file_path),
            source_type=ext,
            content="\n\n".join(sections),
            metadata=metadata,
            created_at=now_iso(),
            updated_at=now_iso(),
            file_hash=file_hash,
        )

    # ──────────────────────────── 文件读取 ────────────────────────────

    @staticmethod
    async def _read_file(file_path: str) -> bytes:
        """异步读取文件二进制内容"""
        def _read() -> bytes:
            with open(file_path, "rb") as f:
                return f.read()
        return await asyncio.to_thread(_read)

    # ──────────────────────────── PDF 解析 ────────────────────────────

    async def parse_pdf(self, raw: bytes) -> list[str]:
        """PDF 解析

        解析策略：
        - 文本型 PDF：pdfplumber 直接提取文本和表格
        - 扫描型 PDF：调用 OCR 微服务识别
        - 混合型：先尝试文本提取，文本量不足的页面走 OCR

        Args:
            raw: PDF 文件二进制内容

        Returns:
            解析后的文本段落列表
        """
        sections: list[str] = []

        try:
            import pdfplumber
        except ImportError:
            logger.error("pdfplumber 未安装，请执行 pip install pdfplumber")
            raise

        def _extract() -> list[str]:
            result: list[str] = []
            with pdfplumber.open(io.BytesIO(raw)) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    tables = page.extract_tables() or []

                    # 文本提取
                    if text and len(text.strip()) > 20:
                        result.append(f"[Page {i + 1}]\n{text.strip()}")

                    # 表格转 Markdown
                    for j, table in enumerate(tables):
                        if table:
                            md_table = self._table_to_markdown(table)
                            result.append(f"[Page {i + 1} Table {j + 1}]\n{md_table}")

                    # 质量检查：文本量不足 → 标记 OCR
                    if not text or len(text.strip()) < 100:
                        result.append(f"[Page {i + 1} OCR_PENDING]")
            return result

        sections = await asyncio.to_thread(_extract)

        # 对标记为 OCR_PENDING 的段落，异步调用 OCR
        ocr_tasks: list[asyncio.Task] = []
        ocr_indices: list[int] = []
        for idx, section in enumerate(sections):
            if "OCR_PENDING" in section:
                page_match = re.search(r"Page (\d+)", section)
                page_num = int(page_match.group(1)) - 1 if page_match else 0
                ocr_indices.append(idx)
                ocr_tasks.append(asyncio.create_task(self._ocr_page(raw, page_num)))

        if ocr_tasks:
            ocr_results = await asyncio.gather(*ocr_tasks, return_exceptions=True)
            for idx, result in zip(ocr_indices, ocr_results):
                if isinstance(result, Exception):
                    logger.warning(f"OCR 失败: {result}")
                    sections[idx] = f"[OCR Failed]"
                else:
                    page_num = idx  # 粗略对应
                    sections[idx] = f"[Page {idx + 1} OCR]\n{result}"

        return sections

    # ──────────────────────────── DOCX 解析 ────────────────────────────

    async def parse_docx(self, raw: bytes) -> list[str]:
        """Word 文档解析

        保留标题层级，提取段落文本和表格，表格转 Markdown。

        Args:
            raw: DOCX 文件二进制内容

        Returns:
            解析后的文本段落列表
        """
        try:
            from docx import Document as DocxDocument
        except ImportError:
            logger.error("python-docx 未安装，请执行 pip install python-docx")
            raise

        def _extract() -> list[str]:
            doc = DocxDocument(io.BytesIO(raw))
            sections: list[str] = []

            for para in doc.paragraphs:
                if para.style and para.style.name and para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "").strip()
                    prefix = "#" * (int(level) if level.isdigit() else 1)
                    sections.append(f"\n{prefix} {para.text}\n")
                elif para.text.strip():
                    sections.append(para.text)

            # 处理表格
            for table in doc.tables:
                md = self._docx_table_to_markdown(table)
                if md:
                    sections.append(md)

            return sections

        return await asyncio.to_thread(_extract)

    # ──────────────────────────── Markdown 解析 ────────────────────────────

    async def parse_markdown(self, raw: bytes) -> list[str]:
        """Markdown 解析

        按标题和代码块分段，保持原始 Markdown 结构。

        Args:
            raw: Markdown 文件二进制内容

        Returns:
            解析后的文本段落列表
        """
        text: str = raw.decode("utf-8", errors="replace")
        sections: list[str] = []

        # 按双换行分段，代码块整体保留
        in_code_block: bool = False
        code_buffer: list[str] = []

        for line in text.split("\n"):
            if line.strip().startswith("```"):
                if in_code_block:
                    code_buffer.append(line)
                    sections.append("\n".join(code_buffer))
                    code_buffer = []
                    in_code_block = False
                else:
                    if code_buffer:
                        sections.append("\n".join(code_buffer))
                        code_buffer = []
                    in_code_block = True
                    code_buffer.append(line)
            elif in_code_block:
                code_buffer.append(line)
            else:
                if line.strip():
                    sections.append(line)

        if code_buffer:
            sections.append("\n".join(code_buffer))

        return sections

    # ──────────────────────────── HTML 解析 ────────────────────────────

    async def parse_html(self, raw: bytes) -> list[str]:
        """HTML 解析

        提取纯文本内容，保留标题层级，表格转 Markdown。

        Args:
            raw: HTML 文件二进制内容

        Returns:
            解析后的文本段落列表
        """
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            logger.error("beautifulsoup4 未安装，请执行 pip install beautifulsoup4")
            raise

        def _extract() -> list[str]:
            text: str = raw.decode("utf-8", errors="replace")
            soup = BeautifulSoup(text, "html.parser")

            # 移除脚本和样式
            for tag in soup(["script", "style", "nav", "footer"]):
                tag.decompose()

            sections: list[str] = []

            # 提取标题
            for level in range(1, 7):
                for heading in soup.find_all(f"h{level}"):
                    prefix = "#" * level
                    sections.append(f"\n{prefix} {heading.get_text().strip()}\n")

            # 提取段落
            for para in soup.find_all("p"):
                text_content = para.get_text().strip()
                if text_content:
                    sections.append(text_content)

            # 提取表格
            for table in soup.find_all("table"):
                md = self._html_table_to_markdown(table)
                if md:
                    sections.append(md)

            return sections

        return await asyncio.to_thread(_extract)

    # ──────────────────────────── 纯文本解析 ────────────────────────────

    async def parse_text(self, raw: bytes) -> list[str]:
        """纯文本解析

        按双换行分段。

        Args:
            raw: 文本文件二进制内容

        Returns:
            解析后的文本段落列表
        """
        text: str = raw.decode("utf-8", errors="replace")
        sections: list[str] = [s.strip() for s in text.split("\n\n") if s.strip()]
        return sections

    # ──────────────────────────── CSV 解析 ────────────────────────────

    async def parse_csv(self, raw: bytes) -> list[str]:
        """CSV 解析

        CSV 整体转为 Markdown 表格。

        Args:
            raw: CSV 文件二进制内容

        Returns:
            包含 Markdown 表格的列表
        """
        import csv

        def _extract() -> list[str]:
            text: str = raw.decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(text))
            rows: list[list[str]] = [list(row) for row in reader]
            if not rows:
                return []
            md = self._table_to_markdown(rows)
            return [md]

        return await asyncio.to_thread(_extract)

    # ──────────────────────────── XLSX 解析 ────────────────────────────

    async def parse_xlsx(self, raw: bytes) -> list[str]:
        """Excel 文档解析

        每个 Sheet 转为一个 Markdown 表格。

        Args:
            raw: XLSX 文件二进制内容

        Returns:
            每个 Sheet 对应的 Markdown 表格列表
        """
        try:
            import openpyxl
        except ImportError:
            logger.error("openpyxl 未安装，请执行 pip install openpyxl")
            raise

        def _extract() -> list[str]:
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            sections: list[str] = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows: list[list[str]] = []
                for row in ws.iter_rows(values_only=True):
                    rows.append([str(cell) if cell is not None else "" for cell in row])
                if rows:
                    md = self._table_to_markdown(rows)
                    sections.append(f"## Sheet: {sheet_name}\n\n{md}")

            wb.close()
            return sections

        return await asyncio.to_thread(_extract)

    # ──────────────────────────── PPTX 解析 ────────────────────────────

    async def parse_pptx(self, raw: bytes) -> list[str]:
        """PowerPoint 文档解析

        每页幻灯片提取文本框内容。

        Args:
            raw: PPTX 文件二进制内容

        Returns:
            每页幻灯片的文本内容列表
        """
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx 未安装，请执行 pip install python-pptx")
            raise

        def _extract() -> list[str]:
            prs = Presentation(io.BytesIO(raw))
            sections: list[str] = []

            for i, slide in enumerate(prs.slides):
                texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts.append(text)
                    # 表格
                    if shape.has_table:
                        table_data: list[list[str]] = []
                        for row in shape.table.rows:
                            table_data.append([cell.text.strip() for cell in row.cells])
                        md = self._table_to_markdown(table_data)
                        texts.append(md)
                if texts:
                    sections.append(f"[Slide {i + 1}]\n" + "\n".join(texts))

            return sections

        return await asyncio.to_thread(_extract)

    # ──────────────────────────── 表格转 Markdown ────────────────────────────

    @staticmethod
    def _table_to_markdown(table: list[list[Any]]) -> str:
        """通用表格转 Markdown

        Args:
            table: 二维列表，第一行为表头

        Returns:
            Markdown 格式的表格字符串
        """
        if not table:
            return ""

        # 确保所有单元格为字符串
        str_table: list[list[str]] = [
            [str(cell) if cell is not None else "" for cell in row]
            for row in table
        ]

        header: str = "| " + " | ".join(str_table[0]) + " |"
        separator: str = "| " + " | ".join("---" for _ in str_table[0]) + " |"
        rows: list[str] = [
            "| " + " | ".join(row) + " |" for row in str_table[1:]
        ]
        return "\n".join([header, separator] + rows)

    @staticmethod
    def _docx_table_to_markdown(table: Any) -> str:
        """DOCX 表格转 Markdown

        Args:
            table: python-docx 的 Table 对象

        Returns:
            Markdown 格式的表格字符串
        """
        rows: list[list[str]] = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        return DocumentParser._table_to_markdown(rows)

    @staticmethod
    def _html_table_to_markdown(table: Any) -> str:
        """HTML 表格转 Markdown

        Args:
            table: BeautifulSoup 的 table 标签对象

        Returns:
            Markdown 格式的表格字符串
        """
        rows: list[list[str]] = []
        for tr in table.find_all("tr"):
            cells = tr.find_all(["td", "th"])
            rows.append([cell.get_text().strip() for cell in cells])
        return DocumentParser._table_to_markdown(rows)

    # ──────────────────────────── OCR ────────────────────────────

    async def _ocr_page(self, raw: bytes, page_number: int) -> str:
        """调用 OCR 微服务识别单页内容

        Args:
            raw: PDF 文件二进制内容
            page_number: 页码（0-indexed）

        Returns:
            OCR 识别的文本
        """
        try:
            # 提取单页 PDF
            single_page_pdf: bytes = await self._extract_single_page(raw, page_number)

            async with aiohttp.ClientSession() as session:
                resp = await session.post(
                    self.OCR_SERVICE_URL,
                    data=single_page_pdf,
                    headers={"Content-Type": "application/pdf"},
                    timeout=aiohttp.ClientTimeout(total=60),
                )
                if resp.status == 200:
                    data = await resp.json()
                    return data.get("text", "")
                else:
                    logger.warning(f"OCR 服务返回 {resp.status}")
                    return ""
        except Exception as e:
            logger.warning(f"OCR 请求失败: {e}")
            return ""

    @staticmethod
    async def _extract_single_page(raw: bytes, page_number: int) -> bytes:
        """从 PDF 中提取单页

        Args:
            raw: PDF 文件二进制内容
            page_number: 页码（0-indexed）

        Returns:
            单页 PDF 的二进制内容
        """
        def _extract() -> bytes:
            try:
                import pdfplumber
                from reportlab.pdfgen import canvas as _  # noqa: F401
            except ImportError:
                return raw  # 降级：发送整个 PDF

            # 使用 pypdf 拆分
            try:
                from pypdf import PdfReader, PdfWriter
                reader = PdfReader(io.BytesIO(raw))
                writer = PdfWriter()
                if page_number < len(reader.pages):
                    writer.add_page(reader.pages[page_number])
                    buf = io.BytesIO()
                    writer.write(buf)
                    return buf.getvalue()
            except Exception:
                pass
            return raw

        return await asyncio.to_thread(_extract)

    # ──────────────────────────── 元数据提取 ────────────────────────────

    @staticmethod
    def _extract_metadata(raw: bytes, ext: str, file_path: str) -> dict[str, Any]:
        """提取文件元数据

        Args:
            raw: 文件二进制内容
            ext: 文件扩展名
            file_path: 文件路径

        Returns:
            元数据字典
        """
        metadata: dict[str, Any] = {
            "format": ext,
            "file_size": len(raw),
            "file_name": os.path.basename(file_path),
        }

        # PDF 页数
        if ext == "pdf":
            try:
                import pdfplumber
                with pdfplumber.open(io.BytesIO(raw)) as pdf:
                    metadata["pages"] = len(pdf.pages)
            except Exception:
                pass

        # DOCX 属性
        if ext == "docx":
            try:
                from docx import Document as DocxDocument
                doc = DocxDocument(io.BytesIO(raw))
                core_props = doc.core_properties
                metadata["author"] = core_props.author or ""
                metadata["title"] = core_props.title or ""
                metadata["created"] = core_props.created.isoformat() if core_props.created else ""
            except Exception:
                pass

        return metadata
